#!/usr/bin/env python3
"""
批量文档预处理器 — Document Preprocessor
=========================================
将 PDF/PPTX 课件批量"翻译"为结构化 Markdown，供后续知识库检索使用。

工作流程:
  1. 扫描输入目录中的所有 PDF/PPTX 文件
  2. 逐个文件：渲染每页为图片 → 调用豆包视觉 API → 得到 Markdown
  3. 支持断点续传（中断后重新运行会跳过已处理的页）
  4. 支持并发处理多个页面（加速大文件）
  5. 最终每个源文件输出一个 .md 文件

用法:
  python doc_preprocessor.py <输入目录> [-o 输出目录] [-w 并发数] [--max-pages N]

示例:
  python doc_preprocessor.py "./课件/" -o "./course_md/" -w 3 -v
"""

import os
import sys
import io
import json
import time
import base64
import hashlib
import argparse
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime

# ============================================================
# 配置
# ============================================================
ARK_API_KEY = os.environ.get(
    "ARK_API_KEY",
    "ark-a73d32ae-9cae-42a7-97bc-d5700f069306-e5ac6"
)
ARK_MODEL = os.environ.get("ARK_MODEL", "ep-20260527110933-btjkj")
BASE_URL = "https://ark.cn-beijing.volces.com/api/v3"

SUPPORTED_EXTS = {".pdf", ".pptx", ".ppt"}

# API 限制
MAX_IMAGE_SIZE_MB = 10       # 单张图片上限
MAX_PAYLOAD_MB = 64          # 请求体上限
MAX_PAGES_PER_REQUEST = 3    # 每次请求最多发几页（避免过大）
DPI = 150                     # PDF 渲染 DPI（平衡速度与清晰度）


# ============================================================
# 文件渲染
# ============================================================

def render_pdf_pages(file_path: str, dpi: int = 150) -> list:
    """将 PDF 每页渲染为 base64 PNG data URI 列表."""
    import fitz
    doc = fitz.open(file_path)
    images = []
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    for page in doc:
        pix = page.get_pixmap(matrix=matrix, colorspace="rgb")
        img_bytes = pix.tobytes("png")
        b64 = base64.b64encode(img_bytes).decode("ascii")
        uri = f"data:image/png;base64,{b64}"
        images.append(uri)
    doc.close()
    return images


def render_pptx_pages(file_path: str) -> list:
    """将 PPTX 每页渲染为 base64 PNG data URI 列表."""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(file_path)
    images = []
    W, H = 1280, 720

    # 尝试加载字体
    font = None
    for fn in ["msyh.ttc", "simhei.ttf", "simsun.ttc", "arial.ttf"]:
        try:
            font = ImageFont.truetype(fn, 20)
            break
        except (IOError, OSError):
            continue

    for slide in prs.slides:
        img = Image.new("RGB", (W, H), "white")
        draw = ImageDraw.Draw(img)
        y = 10

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    chars_per_line = 60
                    while len(text) > 0:
                        chunk = text[:chars_per_line]
                        text = text[chars_per_line:]
                        if font:
                            draw.text((10, y), chunk, fill="black", font=font)
                        else:
                            draw.text((10, y), chunk, fill="black")
                        y += 26
                        if y > H - 10:
                            break
                if y > H - 10:
                    break

            if shape.shape_type == 13:  # PICTURE
                try:
                    image_blob = shape.image.blob
                    pil_img = Image.open(io.BytesIO(image_blob)).convert("RGB")
                    max_w, max_h = 300, 200
                    pil_img.thumbnail((max_w, max_h), Image.LANCZOS)
                    img.paste(pil_img, (W - max_w - 10, y))
                    y += max_h + 10
                except Exception:
                    pass

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        images.append(f"data:image/png;base64,{b64}")

    return images


# ============================================================
# API 调用
# ============================================================

def call_vision_api(image_uris: list, page_numbers: list, source_name: str,
                    max_tokens: int = 4096) -> str:
    """调用豆包视觉 API，识别一页或多页并返回 Markdown."""
    import requests

    prompt = (
        "请将以下课件页面的内容整理为结构化 Markdown 格式。要求：\n"
        "1. 保留所有标题层级（使用 # ## ### 等）\n"
        "2. 保留所有正文文字原文\n"
        "3. 数学公式用 LaTeX 格式（$...$ 或 $$...$$）\n"
        "4. 如有表格，用 Markdown 表格格式\n"
        "5. 如图片中有图表/流程图，用文字描述\n"
        "6. 不要遗漏任何信息\n"
        f"7. 这是文件「{source_name}」的第 {page_numbers} 页"
    )

    content = []
    for uri in image_uris:
        content.append({"type": "image_url", "image_url": {"url": uri}})
    content.append({"type": "text", "text": prompt})

    for attempt in range(3):
        try:
            resp = requests.post(
                f"{BASE_URL}/chat/completions",
                headers={
                    "Authorization": f"Bearer {ARK_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": ARK_MODEL,
                    "messages": [{"role": "user", "content": content}],
                    "max_tokens": max_tokens,
                    "temperature": 0.2,
                },
                timeout=300,
            )
            if resp.status_code == 200:
                return resp.json()["choices"][0]["message"]["content"]
            if resp.status_code == 429:
                wait = min(2 ** attempt * 5, 60)
                print(f"    [RATE-LIMIT] 等待 {wait}s 后重试...", file=sys.stderr)
                time.sleep(wait)
                continue
            raise RuntimeError(f"API 返回 {resp.status_code}: {resp.text[:300]}")
        except requests.exceptions.Timeout:
            print(f"    [TIMEOUT] 重试 {attempt + 1}/3 ...", file=sys.stderr)
            time.sleep(5)
        except requests.exceptions.ConnectionError:
            print(f"    [CONN-ERR] 重试 {attempt + 1}/3 ...", file=sys.stderr)
            time.sleep(10)
    raise RuntimeError("API 调用失败（重试3次后仍失败）")


# ============================================================
# 状态管理（断点续传）
# ============================================================

def load_state(state_file: str) -> dict:
    """加载处理进度."""
    if os.path.exists(state_file):
        with open(state_file, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_state(state_file: str, state: dict):
    """保存处理进度."""
    with open(state_file, "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=2)


# ============================================================
# 主处理流程
# ============================================================

def process_file(file_path: str, output_dir: str, state: dict, workers: int,
                 max_pages: int, dpi: int, verbose: bool) -> dict:
    """处理单个文件，返回更新的 state."""
    fname = Path(file_path).name
    fkey = hashlib.md5(file_path.encode()).hexdigest()[:12]

    if fkey not in state:
        state[fkey] = {"file": file_path, "completed_pages": [], "total_pages": 0}

    fs = state[fkey]
    ext = Path(file_path).suffix.lower()

    # 渲染所有页面
    print(f"\n{'='*60}")
    print(f"[PROCESS] {fname}")
    if verbose:
        print(f"[INFO] 渲染页面中...")

    if ext == ".pdf":
        all_images = render_pdf_pages(file_path, dpi=dpi)
    elif ext in (".pptx", ".ppt"):
        all_images = render_pptx_pages(file_path)
    else:
        print(f"[SKIP] 不支持的文件类型: {ext}")
        return state

    total_pages = min(len(all_images), max_pages)
    fs["total_pages"] = total_pages

    if verbose:
        print(f"[INFO] 共 {total_pages} 页，已完成 {len(fs['completed_pages'])} 页")

    # 确定待处理页面
    pending = [i for i in range(total_pages) if i not in fs["completed_pages"]]
    if not pending:
        print(f"[DONE] 全部已完成，跳过")
        return state

    print(f"[INFO] 待处理: {len(pending)} 页，并发数: {workers}")

    # 并发处理
    def process_page_range(start_idx: int, count: int) -> tuple:
        """处理连续的几页."""
        indices = list(range(start_idx, min(start_idx + count, total_pages)))
        uris = [all_images[i] for i in indices]
        page_nums = [i + 1 for i in indices]  # 1-indexed
        try:
            text = call_vision_api(uris, page_nums, fname)
            return indices, text
        except Exception as e:
            print(f"    [FAIL] 第{page_nums}页: {e}", file=sys.stderr)
            raise

    # 将待处理页面分批，每批 MAX_PAGES_PER_REQUEST 页
    batches = []
    for i in range(0, len(pending), MAX_PAGES_PER_REQUEST):
        batch_pages = pending[i:i + MAX_PAGES_PER_REQUEST]
        batches.append(batch_pages)

    if verbose:
        print(f"[INFO] 共 {len(batches)} 个批次")

    completed_count = 0
    output_file = os.path.join(output_dir, Path(file_path).stem + ".md")

    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {}
        for batch in batches:
            future = executor.submit(process_page_range, batch[0], len(batch))
            futures[future] = batch

        for future in as_completed(futures):
            batch = futures[future]
            try:
                indices, text = future.result()
                # 追加到输出文件
                with open(output_file, "a", encoding="utf-8") as f:
                    f.write(f"\n\n<!-- 第 {indices[0]+1}-{indices[-1]+1} 页 -->\n\n")
                    f.write(text)

                # 更新状态
                for idx in indices:
                    if idx not in fs["completed_pages"]:
                        fs["completed_pages"].append(idx)
                fs["completed_pages"].sort()
                save_state(os.path.join(output_dir, ".preprocessor_state.json"), state)

                completed_count += len(indices)
                pct = len(fs["completed_pages"]) * 100 // total_pages
                print(f"  [{pct:3d}%] 第 {indices[0]+1}-{indices[-1]+1} 页完成 "
                      f"({len(fs['completed_pages'])}/{total_pages})")

            except Exception as e:
                print(f"  [ERROR] 第{batch[0]+1}页起处理失败: {e}", file=sys.stderr)
                # 不退出，继续处理其他批次

    print(f"[DONE] {fname} — {len(fs['completed_pages'])}/{total_pages} 页")
    return state


def main():
    parser = argparse.ArgumentParser(
        description="批量文档预处理器 — PDF/PPTX → Markdown（调用豆包视觉API）"
    )
    parser.add_argument("input_dir", help="输入目录（含 PDF/PPTX 文件）")
    parser.add_argument("-o", "--output-dir", default=None, help="输出目录（默认: 输入目录/md_output）")
    parser.add_argument("-w", "--workers", type=int, default=2, help="并发处理数 (默认2, 建议不超过3)")
    parser.add_argument("--max-pages", type=int, default=500, help="单文件最大处理页数 (默认500)")
    parser.add_argument("--dpi", type=int, default=DPI, help=f"PDF渲染DPI (默认{DPI})")
    parser.add_argument("-v", "--verbose", action="store_true", help="详细输出")
    parser.add_argument("--reset", action="store_true", help="清除进度重新开始")

    args = parser.parse_args()

    input_dir = os.path.abspath(args.input_dir)
    if not os.path.isdir(input_dir):
        print(f"[ERROR] 输入目录不存在: {input_dir}", file=sys.stderr)
        sys.exit(1)

    output_dir = args.output_dir or os.path.join(input_dir, "md_output")
    os.makedirs(output_dir, exist_ok=True)

    state_file = os.path.join(output_dir, ".preprocessor_state.json")
    state = {} if args.reset else load_state(state_file)

    # 扫描文件
    files = []
    for root, _, filenames in os.walk(input_dir):
        # 跳过输出目录
        if os.path.abspath(root).startswith(os.path.abspath(output_dir)):
            continue
        for fname in filenames:
            ext = Path(fname).suffix.lower()
            if ext in SUPPORTED_EXTS:
                files.append(os.path.join(root, fname))

    if not files:
        print(f"[ERROR] 未找到 PDF/PPTX 文件于: {input_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"[INFO] 找到 {len(files)} 个文件:")
    for f in files:
        print(f"  - {Path(f).name}")

    # 逐个处理
    start_time = time.time()
    for file_path in files:
        try:
            state = process_file(
                file_path, output_dir, state,
                args.workers, args.max_pages, args.dpi, args.verbose
            )
        except KeyboardInterrupt:
            print(f"\n[BREAK] 用户中断，进度已保存")
            break

    elapsed = time.time() - start_time
    total_completed = sum(len(s.get("completed_pages", [])) for s in state.values())
    total_pages = sum(s.get("total_pages", 0) for s in state.values())
    print(f"\n{'='*60}")
    print(f"[SUMMARY] 总页数: {total_completed}/{total_pages} | "
          f"耗时: {elapsed/60:.1f}min | 输出: {output_dir}")
    print(f"[SUMMARY] 进度文件: {state_file}")


if __name__ == "__main__":
    main()
