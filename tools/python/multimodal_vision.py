#!/usr/bin/env python3
"""
多模态视觉理解工具 — Multimodal Vision Tool
=============================================
调用豆包 Doubao-Seed-2.0-Mini (火山引擎 Ark API) 进行图片/PDF/PPT的视觉识别。

用法:
    python multimodal_vision.py <file> [-p "提示词"] [--max-pages 5] [-v]

环境变量:
    ARK_API_KEY  — 火山引擎 API Key (可选，脚本内置默认值)
    ARK_MODEL    — 推理接入点 ID (可选，脚本内置默认值)

依赖:
    pip install volcenginesdkarkruntime PyMuPDF python-pptx Pillow
"""

import os
import sys
import io
import base64
import argparse
import json
from pathlib import Path

# ============================================================
# 默认配置 — 可通过环境变量覆盖
# ============================================================
_ARK_API_KEY = "ark-a73d32ae-9cae-42a7-97bc-d5700f069306-e5ac6"
_ARK_MODEL   = "ep-20260528213610-cl26k"  # Doubao-Seed-2.0-lite (2026-05-28 upgraded from Mini)
_BASE_URL    = "https://ark.cn-beijing.volces.com/api/v3"

ARK_API_KEY = os.environ.get("ARK_API_KEY", _ARK_API_KEY)
ARK_MODEL   = os.environ.get("ARK_MODEL", _ARK_MODEL)

# 支持的图片后缀
IMAGE_EXTS = frozenset({".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"})

# MIME 类型映射
_MIME_MAP = {
    "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "png": "image/png", "gif": "image/gif",
    "webp": "image/webp", "bmp": "image/bmp",
}


# ============================================================
# 核心函数
# ============================================================

def file_to_base64_uri(file_path: str) -> str:
    """将图片文件编码为 base64 data URI."""
    ext = Path(file_path).suffix.lower().lstrip(".")
    mime = _MIME_MAP.get(ext, "image/png")
    with open(file_path, "rb") as f:
        data = base64.b64encode(f.read()).decode("ascii")
    return f"data:{mime};base64,{data}"


def pdf_to_base64_uris(file_path: str, dpi: int = 200, max_pages: int = 10) -> list:
    """将 PDF 逐页渲染为 base64 data URI 列表 (需要 PyMuPDF)."""
    import fitz
    doc = fitz.open(file_path)
    images = []
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    for page in doc:
        if len(images) >= max_pages:
            break
        pix = page.get_pixmap(matrix=matrix, colorspace="rgb")
        img_bytes = pix.tobytes("png")
        b64 = base64.b64encode(img_bytes).decode("ascii")
        images.append(f"data:image/png;base64,{b64}")
    doc.close()
    return images


def pptx_to_base64_uris(file_path: str, max_pages: int = 10) -> list:
    """将 PPTX 每页渲染为图片 (需要 python-pptx, Pillow).

    原理：用 PIL 在白色画布上绘制幻灯片中的文本和嵌入图片.
    这并非像素级渲染，但能完整保留文字内容和图片."""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(file_path)
    slide_images = []

    # 尝试加载中文字体
    font = None
    for font_name in ["msyh.ttc", "simhei.ttf", "simsun.ttc", "arial.ttf"]:
        try:
            font = ImageFont.truetype(font_name, 20)
            break
        except (IOError, OSError):
            continue
    # 如果找不到中文字体，用默认字体 (可能不支持中文)
    if font is None:
        try:
            font = ImageFont.load_default()
        except Exception:
            font = None

    # 16:9 画布
    W, H = 1280, 720

    for i, slide in enumerate(prs.slides):
        if len(slide_images) >= max_pages:
            break

        img = Image.new("RGB", (W, H), "white")
        draw = ImageDraw.Draw(img)
        y = 10

        for shape in slide.shapes:
            # 文本提取
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # 近似换行: 每 60 个英文字符或 30 个中文字符换行
                    chars_per_line = 60
                    while len(text) > 0:
                        chunk = text[:chars_per_line]
                        text = text[chars_per_line:]
                        if font:
                            draw.text((10, y), chunk, fill="black", font=font)
                        else:
                            draw.text((10, y), chunk, fill="black")
                        y += 26
                        if y > H - 10:
                            break
                if y > H - 10:
                    break

            # 嵌入图片提取
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_blob = shape.image.blob
                    pil_img = Image.open(io.BytesIO(image_blob))
                    pil_img = pil_img.convert("RGB")
                    # 缩放到合适大小
                    max_w, max_h = 300, 200
                    pil_img.thumbnail((max_w, max_h), Image.LANCZOS)
                    img.paste(pil_img, (W - max_w - 10, y))
                    y += max_h + 10
                except Exception:
                    pass

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        slide_images.append(f"data:image/png;base64,{b64}")

    return slide_images


def pptx_extract_text(file_path: str, max_pages: int = 10) -> str:
    """从 PPTX 提取纯文本 (备用方案: 不需要视觉能力)."""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides):
        if len(parts) >= max_pages:
            break
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            parts.append(f"【第 {i + 1} 页】\n" + "\n".join(texts))
    return "\n\n".join(parts)


def call_chat_api(messages: list, max_tokens: int = 4096, temperature: float = 0.3) -> str:
    """调用火山引擎 Chat Completions API."""
    import requests

    payload = {
        "model": ARK_MODEL,
        "messages": messages,
        "max_tokens": max_tokens,
        "temperature": temperature,
        "thinking": {"type": "disabled"},  # 纯视觉识别不需要深度思考，降延迟
    }
    headers = {
        "Authorization": f"Bearer {ARK_API_KEY}",
        "Content-Type": "application/json",
    }

    resp = requests.post(
        f"{_BASE_URL}/chat/completions",
        headers=headers,
        json=payload,
        timeout=180,
    )

    if resp.status_code != 200:
        detail = resp.text[:500]
        raise RuntimeError(f"API 返回 {resp.status_code}: {detail}")

    data = resp.json()
    return data["choices"][0]["message"]["content"]


def build_vision_messages(image_uris: list, prompt: str, max_images_per_request: int = 10):
    """构建多模态消息列表.

    豆包 API 支持一个 content 数组包含多个 image_url.
    如果图片太多，分批放在多条 user 消息中."""
    messages = []
    for chunk_start in range(0, len(image_uris), max_images_per_request):
        chunk = image_uris[chunk_start:chunk_start + max_images_per_request]
        content = []
        for uri in chunk:
            content.append({"type": "image_url", "image_url": {"url": uri, "detail": "high"}})
        content.append({"type": "text", "text": prompt})
        messages.append({"role": "user", "content": content})
        prompt = "请继续上述内容，识别下一页/下一批图片。"  # 后续的消息用简短提示
    return messages


# ============================================================
# 命令行入口
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="多模态视觉理解工具 — 调用豆包 API 识别图片/PDF/PPT",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python multimodal_vision.py photo.jpg
  python multimodal_vision.py homework.pdf -p "提取所有数学公式"
  python multimodal_vision.py slides.pptx -p "总结内容" --max-pages 20
  python multimodal_vision.py diagram.png -p "描述这个流程图" -v
        """,
    )
    parser.add_argument("file", help="要识别的文件路径 (支持 jpg/png/gif/webp/pdf/pptx)")
    parser.add_argument("-p", "--prompt", default="请详细、完整地描述该文件的内容。包括：所有文字（保留原文）、公式（LaTeX格式）、表格数据、图表内容、图像描述等。不要遗漏任何信息。", help="自定义提示词")
    parser.add_argument("--max-pages", type=int, default=10, help="PDF/PPTX 最多处理页数 (默认10)")
    parser.add_argument("--dpi", type=int, default=200, help="PDF 渲染分辨率 (默认200)")
    parser.add_argument("--text-only", action="store_true", help="PPTX仅提取文本，不渲染图像")
    parser.add_argument("-v", "--verbose", action="store_true", help="显示处理进度信息")
    parser.add_argument("--no-vision", action="store_true", help="强制纯文本模式 (不发送图片)")
    parser.add_argument("-o", "--output", help="将识别结果写入指定文件 (默认输出到 stdout)")

    args = parser.parse_args()
    file_path = args.file
    silent = lambda *a, **kw: None

    def output(text: str):
        """输出结果到 stdout 或文件."""
        if args.output:
            with open(args.output, "w", encoding="utf-8") as f:
                f.write(text)
            print(f"[INFO] 结果已写入: {args.output}", file=sys.stderr)
        else:
            # 确保 stdout 使用 UTF-8
            sys.stdout.reconfigure(encoding="utf-8")
            print(text)

    if not os.path.exists(file_path):
        print(f"[ERROR] 文件不存在: {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = Path(file_path).suffix.lower()
    file_size = os.path.getsize(file_path)
    log = print if args.verbose else silent

    log(f"[INFO] 文件: {os.path.basename(file_path)}")
    log(f"[INFO] 类型: {ext}  |  大小: {file_size / 1024:.1f} KB")
    log(f"[INFO] 模型: {ARK_MODEL}")

    image_uris = []

    # ---- 分类处理 ----
    if ext in IMAGE_EXTS:
        log("[INFO] 图片文件，直接编码...")
        image_uris = [file_to_base64_uri(file_path)]
        log(f"[INFO] base64 编码完成 ({len(image_uris[0]):,} 字符)")

    elif ext == ".pdf":
        log(f"[INFO] PDF文件，渲染中 (DPI={args.dpi})...")
        try:
            image_uris = pdf_to_base64_uris(file_path, dpi=args.dpi, max_pages=args.max_pages)
            log(f"[INFO] 已渲染 {len(image_uris)} 页")
        except ImportError:
            print("[ERROR] 处理PDF需要 PyMuPDF: pip install PyMuPDF", file=sys.stderr)
            sys.exit(2)

    elif ext in (".ppt", ".pptx"):
        if args.no_vision or args.text_only:
            log("[INFO] PPTX纯文本模式...")
            try:
                text = pptx_extract_text(file_path, max_pages=args.max_pages)
            except ImportError:
                print("[ERROR] 需要 python-pptx: pip install python-pptx", file=sys.stderr)
                sys.exit(2)
            # 直接用文本调用非视觉API
            log("[INFO] 发送文本到API...")
            try:
                result = call_chat_api([
                    {"role": "user", "content": f"{args.prompt}\n\n文件内容如下:\n\n{text}"}
                ])
                output(result)
                return
            except Exception as e:
                print(f"[ERROR] API调用失败: {e}", file=sys.stderr)
                sys.exit(3)
        else:
            log("[INFO] PPTX文件，渲染幻灯片...")
            try:
                image_uris = pptx_to_base64_uris(file_path, max_pages=args.max_pages)
                log(f"[INFO] 已渲染 {len(image_uris)} 页")
            except ImportError as e:
                print(f"[ERROR] 处理PPTX需要 python-pptx 和 Pillow: pip install python-pptx Pillow", file=sys.stderr)
                sys.exit(2)
    else:
        print(f"[ERROR] 不支持的文件类型: {ext}", file=sys.stderr)
        print(f"[INFO] 支持的格式: {', '.join(sorted(IMAGE_EXTS))}, .pdf, .pptx", file=sys.stderr)
        sys.exit(1)

    if not image_uris:
        print("[ERROR] 未能从文件中提取任何图像", file=sys.stderr)
        sys.exit(1)

    # ---- 调用视觉API ----
    log(f"[INFO] 共 {len(image_uris)} 张图片，准备调用视觉API...")

    try:
        msgs = build_vision_messages(image_uris, args.prompt)
        if len(msgs) > 1 and not args.verbose:
            # 如果有多条消息，每批单独请求并合并结果
            all_results = []
            for i, msg in enumerate(msgs):
                log(f"[INFO] 批次 {i + 1}/{len(msgs)} ...")
                content = call_chat_api([msg])
                all_results.append(content)
            output("\n\n---\n\n".join(all_results))
        else:
            result = call_chat_api(msgs)
            output(result)
    except Exception as e:
        print(f"[ERROR] API调用失败: {e}", file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()
